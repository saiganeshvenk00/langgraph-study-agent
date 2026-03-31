from pathlib import Path
import pypdf
from docx import Document
from pptx import Presentation


#this is the default notes folder - it points to the Notes folder in the project directory
NOTES_DIR = Path("Notes")

#these are all the file types the agent can read from the notes folder
SUPPORTED_EXTENSIONS = ["*.txt", "*.pdf", "*.docx", "*.pptx"]


#this definition is used to update the notes folder path at runtime when the user sets a new one from the UI
def set_notes_dir(path: str):
    global NOTES_DIR
    NOTES_DIR = Path(path)


#this definition is used to list all the note files in the Notes directory across all supported formats
def list_note_files() -> list[str]:
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        files.extend([f.name for f in NOTES_DIR.glob(ext)])
    return files


#this definition is used to read the content of a note file - it detects the file type and parses it accordingly
def read_note_file(file_name: str) -> str:
    file_path = NOTES_DIR / file_name
    suffix = file_path.suffix.lower()

    #if its a plain text file, just read it directly
    if suffix == ".txt":
        return file_path.read_text(encoding="utf-8")

    #if its a pdf, loop through all the pages and extract the text from each one
    elif suffix == ".pdf":
        reader = pypdf.PdfReader(str(file_path))
        return "\n".join(
            page.extract_text() for page in reader.pages if page.extract_text()
        )

    #if its a word doc, loop through all the paragraphs and grab the text
    elif suffix == ".docx":
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    #if its a powerpoint, loop through every slide, then every shape, then every paragraph to get the text
    elif suffix == ".pptx":
        prs = Presentation(str(file_path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text.append(para.text)
        return "\n".join(text)

    #if the file type is not supported, return a message saying so instead of crashing
    else:
        return f"Unsupported file format: {suffix}"


#this definition is used to save the output to a file
def save_output(content: str, filename: str = "output.txt") -> str:
    output_path = Path(filename)
    output_path.write_text(content, encoding="utf-8")
    return str(output_path)
